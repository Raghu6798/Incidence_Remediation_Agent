from pptx import Presentation
import os


def extract_images_from_pptx(pptx_path, output_dir):
    # Load the presentation
    prs = Presentation(pptx_path)

    # Create the output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    image_count = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # shape_type 13 corresponds to a Picture
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext  # 'jpeg', 'png', etc.

                image_filename = f"image_{image_count}.{image_format}"
                image_path = os.path.join(output_dir, image_filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                print(f"Saved: {image_path}")
                image_count += 1

    if image_count == 0:
        print("No images found in the presentation.")
    else:
        print(f"Extracted {image_count} images to: {output_dir}")


# Example usage
extract_images_from_pptx("/content/Test Case HackRx.pptx", "extracted_images")
